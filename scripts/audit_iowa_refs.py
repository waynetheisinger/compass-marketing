"""
Find every text reference to 'IOWA'/'Iowa'/'uiowa' across the template's
slide master, layouts, and existing sample slides — so we know what needs
stripping before we add MowDirect content.
"""
import re
from pptx import Presentation

TEMPLATE = "/Users/waynetheisinger/Desktop/PPT-Template-Standard-2025.pptx"
PATTERN = re.compile(r"iowa", re.IGNORECASE)

prs = Presentation(TEMPLATE)


def scan_shape(shape, location):
    hits = []
    if shape.has_text_frame:
        for pi, para in enumerate(shape.text_frame.paragraphs):
            for ri, run in enumerate(para.runs):
                if PATTERN.search(run.text or ""):
                    hits.append((location, f"para{pi}.run{ri}", shape.name, run.text))
    if shape.shape_type == 6:  # group
        for sub in shape.shapes:
            hits.extend(scan_shape(sub, location))
    return hits


# Slide master
print("=== SLIDE MASTER ===")
for shape in prs.slide_master.shapes:
    for h in scan_shape(shape, "master"):
        print(f"  {h}")

# Layouts
print("\n=== LAYOUTS ===")
for li, layout in enumerate(prs.slide_masters[0].slide_layouts):
    for shape in layout.shapes:
        for h in scan_shape(shape, f"layout[{li}]:{layout.name}"):
            print(f"  {h}")

# Slides
print("\n=== SLIDES ===")
for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        for h in scan_shape(shape, f"slide[{si}]"):
            print(f"  {h}")

# Also check picture shapes for IOWA-themed names
print("\n=== PICTURE SHAPES (name audit) ===")
for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13 and "iowa" in (shape.name or "").lower():  # picture
            print(f"  slide[{si}] picture name: {shape.name}")

for layout in prs.slide_masters[0].slide_layouts:
    for shape in layout.shapes:
        if shape.shape_type == 13 and "iowa" in (shape.name or "").lower():
            print(f"  layout {layout.name!r} picture name: {shape.name}")

for shape in prs.slide_master.shapes:
    if shape.shape_type == 13 and "iowa" in (shape.name or "").lower():
        print(f"  master picture name: {shape.name}")
