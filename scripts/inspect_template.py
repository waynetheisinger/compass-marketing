"""
Quick inspector for the PPT template — lists slide masters, layouts, and any
existing slides so we know what we have to work with before building the deck.
"""
from pptx import Presentation
from pptx.util import Emu

TEMPLATE = "/Users/waynetheisinger/Desktop/PPT-Template-Standard-2025.pptx"

prs = Presentation(TEMPLATE)

print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU")
print(f"           = {Emu(prs.slide_width).inches:.2f}\" x {Emu(prs.slide_height).inches:.2f}\"")
print()

print(f"Slide masters: {len(prs.slide_masters)}")
for mi, master in enumerate(prs.slide_masters):
    print(f"  Master {mi}: {master.name!r}")
    print(f"    Layouts ({len(master.slide_layouts)}):")
    for li, layout in enumerate(master.slide_layouts):
        ph_names = [(ph.placeholder_format.idx, ph.placeholder_format.type, ph.name)
                    for ph in layout.placeholders]
        print(f"      [{li}] {layout.name!r}  placeholders={ph_names}")
print()

print(f"Existing slides: {len(prs.slides)}")
for si, slide in enumerate(prs.slides):
    layout = slide.slide_layout.name
    text_summary = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    text_summary.append(t)
    summary = " | ".join(text_summary[:3])[:120]
    print(f"  Slide {si}: layout={layout!r}  text={summary!r}")
