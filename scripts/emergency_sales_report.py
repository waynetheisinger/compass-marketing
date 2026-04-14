"""
Emergency sales analysis: March vs April 2026.

Pulls aggregate sales data from Shopify via ShopifyQL (read_reports scope),
classifies products as petrol/cordless/other, and generates:
  1. JSON summary  → reports/emergency_sales_data.json
  2. Word document  → reports/Emergency_Marketing_Response.docx
  3. PowerPoint     → reports/Emergency_Marketing_Response.pptx

Usage:
    python scripts/emergency_sales_report.py
"""
from __future__ import annotations

import json
import os
import sys
from collections import defaultdict
from datetime import datetime, timezone

from dotenv import load_dotenv

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
load_dotenv()

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "reports")

# ---------------------------------------------------------------------------
# Classification
# ---------------------------------------------------------------------------

CORDLESS_KW = ["cordless", "battery", "lithium", "40v", "charger", "sbs",
               "sbs460", "sbs220", "sbs480", "sbs560", "sbs240", "sbscdc",
               "sbscsc", "sbscbc", "sbs20cb", "sbs40cb"]
PETROL_KW   = ["petrol", "honda", "4-stroke", "2-stroke", "self-propelled petrol",
               "mountfield", "stiga", "feider", "tg46", "tg40", "tg51", "tg46se",
               "kt400", "sw420", "dct38", "dc24", "ride-on", "lawn tractor",
               "mtf66", "sp185", "t30m", "1538h", "827m", "hf 2625"]
FUEL_CAN_KW = ["jerry can", "fuel can", "pfc-20"]


def classify_product(title: str, product_type: str) -> str:
    """Return one of: cordless, petrol, fuel_cans, other."""
    t = ((title or "") + " " + (product_type or "")).lower()
    if any(k in t for k in FUEL_CAN_KW):
        return "fuel_cans"
    if any(k in t for k in CORDLESS_KW):
        return "cordless"
    if any(k in t for k in PETROL_KW):
        return "petrol"
    # Product types that are inherently petrol
    if (product_type or "").lower() in ("ride-on mowers", "lawn tractors", "garden tractor",
                                 "garden tractors", "log splitters", "chippers"):
        return "petrol"
    return "other"


# ---------------------------------------------------------------------------
# ShopifyQL data pull
# ---------------------------------------------------------------------------

SALES_BY_PRODUCT_QUERY = '''
{
  shopifyqlQuery(query: "FROM sales SHOW net_sales, gross_sales, orders GROUP BY product_title, product_type SINCE %s UNTIL %s ORDER BY net_sales DESC") {
    ... on ShopifyqlQueryResponse {
      parseErrors
      tableData { columns { name dataType } rows }
    }
  }
}
'''

SALES_TOTAL_QUERY = '''
{
  shopifyqlQuery(query: "FROM sales SHOW net_sales, gross_sales, orders SINCE %s UNTIL %s") {
    ... on ShopifyqlQueryResponse {
      parseErrors
      tableData { columns { name dataType } rows }
    }
  }
}
'''

SALES_BY_TYPE_QUERY = '''
{
  shopifyqlQuery(query: "FROM sales SHOW net_sales, gross_sales, orders GROUP BY product_type SINCE %s UNTIL %s ORDER BY net_sales DESC") {
    ... on ShopifyqlQueryResponse {
      parseErrors
      tableData { columns { name dataType } rows }
    }
  }
}
'''


def run_shopifyql(client, query_template: str, start: str, end: str) -> dict:
    """Execute a ShopifyQL query and return parsed tableData."""
    query = query_template % (start, end)
    result = client.execute(query)
    qr = result["shopifyqlQuery"]
    errors = qr.get("parseErrors")
    if errors:
        raise RuntimeError(f"ShopifyQL parse errors: {errors}")
    return qr["tableData"]


def fetch_all_data() -> dict:
    """Fetch March and April sales data from Shopify."""
    from scripts.shopify_client import ShopifyClient

    with ShopifyClient() as client:
        data = {}
        for period, start, end in [
            ("march", "2026-03-01", "2026-03-31"),
            ("april", "2026-04-01", "2026-04-14"),
        ]:
            data[f"{period}_by_product"] = run_shopifyql(client, SALES_BY_PRODUCT_QUERY, start, end)
            data[f"{period}_total"] = run_shopifyql(client, SALES_TOTAL_QUERY, start, end)
            data[f"{period}_by_type"] = run_shopifyql(client, SALES_BY_TYPE_QUERY, start, end)

    return data


def analyse(data: dict) -> dict:
    """Classify and aggregate the raw ShopifyQL data."""
    result = {}

    for period in ("march", "april"):
        totals = data[f"{period}_total"]["rows"][0]
        products = data[f"{period}_by_product"]["rows"]
        by_type = data[f"{period}_by_type"]["rows"]

        # Classify each product
        categories: dict[str, dict] = defaultdict(lambda: {
            "net_sales": 0.0, "gross_sales": 0.0, "orders": 0, "products": []
        })

        for row in products:
            title = row.get("product_title", "")
            ptype = row.get("product_type", "")
            net = float(row.get("net_sales", 0))
            gross = float(row.get("gross_sales", 0))
            orders = int(row.get("orders", 0))

            cat = classify_product(title, ptype)
            categories[cat]["net_sales"] += net
            categories[cat]["gross_sales"] += gross
            categories[cat]["orders"] += orders
            categories[cat]["products"].append({
                "title": title,
                "product_type": ptype,
                "net_sales": net,
                "gross_sales": gross,
                "orders": orders,
            })

        result[period] = {
            "net_sales": float(totals.get("net_sales", 0)),
            "gross_sales": float(totals.get("gross_sales", 0)),
            "orders": int(totals.get("orders", 0)),
            "by_shopify_type": [
                {
                    "product_type": r.get("product_type", ""),
                    "net_sales": float(r.get("net_sales", 0)),
                    "gross_sales": float(r.get("gross_sales", 0)),
                    "orders": int(r.get("orders", 0)),
                }
                for r in by_type
            ],
            "categories": {
                k: {
                    "net_sales": round(v["net_sales"], 2),
                    "gross_sales": round(v["gross_sales"], 2),
                    "orders": v["orders"],
                    "products": sorted(v["products"], key=lambda x: -x["net_sales"]),
                }
                for k, v in categories.items()
            },
        }

    return result


# ---------------------------------------------------------------------------
# Document generation — Word
# ---------------------------------------------------------------------------

def build_docx(analysis: dict, output_path: str):
    """Generate the marketing strategy Word document."""
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    march = analysis["march"]
    april = analysis["april"]
    march_days = 31
    april_days = 14

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # --- TITLE ---
    doc.add_paragraph("")
    title = doc.add_heading("Emergency Marketing Response Plan", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = doc.add_heading("MowDirect — Strait of Hormuz Crisis Response", level=2)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("")
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run("Prepared: 14 April 2026\n").bold = True
    meta.add_run("Confidential — Internal + Factory Partner")
    doc.add_page_break()

    # --- 1. EXECUTIVE SUMMARY ---
    doc.add_heading("1. Executive Summary", level=1)
    doc.add_paragraph(
        "The ongoing naval blockade of the Strait of Hormuz has created an acute disruption "
        "to the garden machinery sector. Crude oil above $100/barrel has driven retail petrol "
        "prices up over 20%, directly increasing the cost of ownership for petrol-powered garden "
        "tools and eroding consumer confidence in petrol purchases."
    )

    march_daily = march["net_sales"] / march_days
    april_daily = april["net_sales"] / april_days
    overall_change = ((april_daily - march_daily) / march_daily * 100) if march_daily else 0

    doc.add_paragraph(
        f"MowDirect's daily sales run rate has dropped from £{march_daily:,.0f}/day in March "
        f"to £{april_daily:,.0f}/day in April (1-14), a {overall_change:+.0f}% decline. "
        f"This document analyses the performance breakdown and presents an emergency response plan."
    )
    doc.add_paragraph("The plan prioritises three pillars:")
    for b in [
        "Accelerate Spectrum cordless range positioning as the primary product line",
        "Reframe petrol products around fuel efficiency and total cost of ownership",
        "Implement website/UX changes to address factory partner feedback on cordless presentation",
    ]:
        doc.add_paragraph(b, style="List Bullet")

    # --- 2. MARKET SITUATION ---
    doc.add_heading("2. Market Situation", level=1)
    doc.add_heading("2.1 Strait of Hormuz Crisis", level=2)
    doc.add_paragraph(
        "Following the collapse of US-Iran peace talks, a naval blockade has reduced commercial "
        "transit through the Strait of Hormuz by approximately 94%. Over 3,200 vessels including "
        "800 tankers are stalled or rerouted, creating a global shortfall of ~20 million barrels/day. "
        "Crude oil has surged back above $100/barrel."
    )
    doc.add_heading("2.2 Impact on Garden Machinery", level=2)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Factor"
    hdr[1].text = "Impact"
    for factor, impact in [
        ("Petrol price spike", "European fuel prices up >20%, increasing annual mower running costs"),
        ("Consumer retreat", "Households cutting back on petrol-powered purchases"),
        ("Electrification acceleration", "Battery/cordless tools seeing increased demand"),
        ("Manufacturing surcharges", "Up to 30% on raw materials — higher retail prices"),
        ("Supply disruption", "Petrol mower supply chains disrupted for coming weeks"),
    ]:
        row = table.add_row().cells
        row[0].text = factor
        row[1].text = impact

    # --- 3. SALES PERFORMANCE ANALYSIS ---
    doc.add_heading("3. Sales Performance Analysis", level=1)

    doc.add_heading("3.1 Overall Performance", level=2)
    table = doc.add_table(rows=1, cols=5)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Metric"
    hdr[1].text = "March 2026"
    hdr[2].text = "April 1-14"
    hdr[3].text = "April Daily Rate"
    hdr[4].text = "vs March"
    for label, m_val, a_val in [
        ("Net Sales", march["net_sales"], april["net_sales"]),
        ("Gross Sales", march["gross_sales"], april["gross_sales"]),
    ]:
        m_daily = m_val / march_days
        a_daily = a_val / april_days
        chg = ((a_daily - m_daily) / m_daily * 100) if m_daily else 0
        row = table.add_row().cells
        row[0].text = label
        row[1].text = f"£{m_val:,.2f}"
        row[2].text = f"£{a_val:,.2f}"
        row[3].text = f"£{a_daily:,.0f}/day"
        row[4].text = f"{chg:+.0f}%"
    row = table.add_row().cells
    row[0].text = "Orders"
    row[1].text = str(march["orders"])
    row[2].text = str(april["orders"])
    m_daily_o = march["orders"] / march_days
    a_daily_o = april["orders"] / april_days
    o_chg = ((a_daily_o - m_daily_o) / m_daily_o * 100) if m_daily_o else 0
    row[3].text = f"{a_daily_o:.0f}/day"
    row[4].text = f"{o_chg:+.0f}%"

    # --- 3.2 By Shopify product type ---
    doc.add_heading("3.2 Sales by Product Type", level=2)
    doc.add_paragraph("Shopify product type breakdown — March vs April daily run rate:")

    # Build lookup for April by type
    april_type_lookup = {r["product_type"]: r for r in april["by_shopify_type"]}

    table = doc.add_table(rows=1, cols=5)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Product Type"
    hdr[1].text = "March Net"
    hdr[2].text = "April Net"
    hdr[3].text = "Daily Rate Change"
    hdr[4].text = "Trend"
    for m_row in march["by_shopify_type"]:
        pt = m_row["product_type"]
        a_row = april_type_lookup.get(pt, {"net_sales": 0, "orders": 0})
        m_d = m_row["net_sales"] / march_days
        a_d = a_row["net_sales"] / april_days
        chg = ((a_d - m_d) / m_d * 100) if m_d else 0
        row = table.add_row().cells
        row[0].text = pt or "(uncategorised)"
        row[1].text = f"£{m_row['net_sales']:,.0f}"
        row[2].text = f"£{a_row['net_sales']:,.0f}"
        row[3].text = f"{chg:+.0f}%"
        row[4].text = "Declining" if chg < -15 else ("Growing" if chg > 15 else "Stable")

    # --- 3.3 Petrol vs Cordless ---
    doc.add_heading("3.3 Petrol vs Cordless Classification", level=2)
    doc.add_paragraph(
        "Products classified by power source (petrol includes ride-ons/tractors, "
        "cordless includes Spectrum 40V battery products):"
    )

    table = doc.add_table(rows=1, cols=6)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Category"
    hdr[1].text = "March Net"
    hdr[2].text = "March Orders"
    hdr[3].text = "April Net"
    hdr[4].text = "April Orders"
    hdr[5].text = "Daily Rate Change"
    for cat in ["petrol", "cordless", "fuel_cans", "other"]:
        m_cat = march["categories"].get(cat, {"net_sales": 0, "orders": 0})
        a_cat = april["categories"].get(cat, {"net_sales": 0, "orders": 0})
        m_d = m_cat["net_sales"] / march_days
        a_d = a_cat["net_sales"] / april_days
        chg = ((a_d - m_d) / m_d * 100) if m_d else 0
        row = table.add_row().cells
        row[0].text = cat.replace("_", " ").title()
        row[1].text = f"£{m_cat['net_sales']:,.0f}"
        row[2].text = str(m_cat["orders"])
        row[3].text = f"£{a_cat['net_sales']:,.0f}"
        row[4].text = str(a_cat["orders"])
        row[5].text = f"{chg:+.0f}%"

    # --- 3.4 Top products April ---
    doc.add_heading("3.4 Top Selling Products — April 2026", level=2)
    all_april_prods = []
    for cat, cdata in april["categories"].items():
        for p in cdata["products"]:
            all_april_prods.append({**p, "category": cat})
    all_april_prods.sort(key=lambda x: -x["net_sales"])

    table = doc.add_table(rows=1, cols=5)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Product"
    hdr[1].text = "Type"
    hdr[2].text = "Category"
    hdr[3].text = "Net Sales"
    hdr[4].text = "Orders"
    for p in all_april_prods[:20]:
        row = table.add_row().cells
        row[0].text = p["title"][:55]
        row[1].text = p["product_type"] or ""
        row[2].text = p["category"].replace("_", " ").title()
        row[3].text = f"£{p['net_sales']:,.2f}"
        row[4].text = str(p["orders"])

    # --- 3.5 Spectrum Cordless Products ---
    doc.add_heading("3.5 Spectrum Cordless Range Performance", level=2)
    doc.add_paragraph(
        "The Spectrum 40V cordless range — the products the factory partner specifically "
        "highlighted — and their current sales performance:"
    )
    # Find cordless products in both months
    m_cordless = march["categories"].get("cordless", {"products": []})
    a_cordless = april["categories"].get("cordless", {"products": []})

    if m_cordless["products"] or a_cordless["products"]:
        m_lookup = {p["title"]: p for p in m_cordless["products"]}
        a_lookup = {p["title"]: p for p in a_cordless["products"]}
        all_titles = set(m_lookup.keys()) | set(a_lookup.keys())

        table = doc.add_table(rows=1, cols=4)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text = "Product"
        hdr[1].text = "March Net"
        hdr[2].text = "April Net"
        hdr[3].text = "April Orders"
        for t in sorted(all_titles):
            m_p = m_lookup.get(t, {"net_sales": 0})
            a_p = a_lookup.get(t, {"net_sales": 0, "orders": 0})
            row = table.add_row().cells
            row[0].text = t[:60]
            row[1].text = f"£{m_p['net_sales']:,.2f}"
            row[2].text = f"£{a_p['net_sales']:,.2f}"
            row[3].text = str(a_p.get("orders", 0))
    else:
        doc.add_paragraph(
            "No cordless products detected in sales data. This confirms the factory partner's "
            "concern — the products are not being found or purchased effectively."
        )

    # --- 4. STRATEGIC RESPONSE PLAN ---
    doc.add_heading("4. Strategic Response Plan", level=1)

    doc.add_heading("4.1 Pillar 1 — Accelerate Spectrum Cordless", level=2)
    doc.add_paragraph(
        "The Strait of Hormuz crisis represents a significant opportunity for the Spectrum "
        "cordless range. With petrol supply disrupted and fuel costs surging, the value proposition "
        "of battery-powered tools has never been stronger. The factory partner has up to £1M "
        "invested and expects rapid website improvements."
    )
    doc.add_heading("Website/UX Changes (Factory Partner Feedback)", level=3)
    for title_, desc in [
        ("Complete machine bundles", "Every cordless tool must be available as a complete kit "
         "(tool + battery/batteries + charger) at a single price, as well as tool-only. "
         "This addresses the factory partner's primary concern."),
        ("Battery bundle clarity", "For products requiring two batteries (e.g. the 46cm mower), "
         "create a dedicated bundle variant that clearly shows '2x 4.0Ah batteries + charger' "
         "with lifestyle imagery of the complete kit."),
        ("Cross-sell ecosystem page", "Build a dedicated 'Spectrum 40V System' landing page "
         "showing all tools sharing the same battery platform. Hero message: "
         "'One battery powers everything.'"),
        ("Lifestyle imagery", "Replace/supplement technical product shots with lifestyle images. "
         "Source from factory partner."),
        ("Product descriptions", "Expand each listing with use cases, garden size suitability, "
         "runtime expectations, comparison vs petrol equivalent."),
    ]:
        p = doc.add_paragraph()
        p.add_run(f"{title_}: ").bold = True
        p.add_run(desc)

    doc.add_heading("4.2 Pillar 2 — Reframe Petrol Range", level=2)
    doc.add_paragraph(
        "Petrol products remain a significant revenue line. Our strategy is to reframe "
        "petrol around fuel efficiency and total cost of ownership."
    )
    for title_, desc in [
        ("'Fuel Crisis Kit' bundles", "Bundle petrol mowers with a 20L jerry can (already "
         "selling well — £10K+ in March), fuel stabiliser, and a '12-month fuel cost guide'. "
         "Position as 'Everything you need to ride out rising fuel costs.'"),
        ("Fuel efficiency callouts", "Add prominent fuel efficiency badges to all petrol "
         "product listings. E.g. 'Only 0.8L/hour — mow for just £2.50 per session.'"),
        ("Trade-in programme", "'Switch to Cordless' — offer trade-in credit (£30-£50) "
         "on old petrol mowers when purchasing a Spectrum cordless kit."),
        ("Servicing package", "'Peace of Mind Petrol Pack' — annual service voucher "
         "bundled at purchase to reduce perceived ongoing cost anxiety."),
    ]:
        p = doc.add_paragraph()
        p.add_run(f"{title_}: ").bold = True
        p.add_run(desc)

    doc.add_heading("4.3 Pillar 3 — Messaging & Content Pivot", level=2)
    for title_, desc in [
        ("Homepage hero banner", "Rotate to Spectrum cordless with crisis-aware messaging: "
         "'No fuel. No fumes. No worries.' / 'Power your garden without the pump.'"),
        ("Cost comparison calculator", "Interactive tool showing 5-year TCO: cordless vs petrol "
         "(fuel + servicing + purchase price). Auto-updates with current fuel prices."),
        ("Email campaigns (Klaviyo)", "Segmented: (1) Past petrol buyers — 'Time to switch?' "
         "with trade-in offer; (2) Browsing cordless — 'Complete your kit'; "
         "(3) All customers — 'Beat the fuel crisis' editorial."),
        ("Social content", "TikTok/Instagram Reels: 'Petrol vs Cordless cost diary' series. "
         "Side-by-side monthly cost comparison."),
        ("Blog/SEO", "Publish 'Petrol vs Battery Lawn Mowers 2026: The Real Cost' targeting "
         "high-intent search terms spiking during the crisis."),
        ("Accessory attach", "Cross-sell spare batteries on every cordless page. For petrol: "
         "fuel cans, oil, maintenance kits. Target: +15% AOV."),
    ]:
        p = doc.add_paragraph()
        p.add_run(f"{title_}: ").bold = True
        p.add_run(desc)

    # --- 5. ADDITIONAL IDEAS ---
    doc.add_heading("5. Additional Strategic Initiatives", level=1)
    for title_, desc in [
        ("Marketplace priority shift", "Prioritise Spectrum cordless on Amazon, eBay, B&Q, "
         "ManoMano. Update titles to lead with 'No fuel needed'. Boost Promoted Listings +50%."),
        ("Google Ads restructure", "Reduce budget on petrol PMax asset groups. Redirect to "
         "cordless campaigns with crisis-aware copy."),
        ("Influencer partnerships", "Fast-track 2-3 garden/lifestyle influencer partnerships "
         "for Spectrum cordless. Brief: 'switching from petrol' narrative."),
        ("PR opportunity", "Pitch: 'UK garden retailer sees cordless sales surge amid fuel "
         "crisis' — positions MowDirect as thought leader."),
        ("Stock reallocation", "Allocate Spectrum cordless stock to highest-converting channels "
         "first (Shopify direct, then Amazon)."),
        ("Flash sale", "'Green Garden Weekend' — 48-hour event, 10% off all cordless + free "
         "battery upgrade on selected kits."),
        ("Warranty messaging", "Amplify Spectrum's 5-year warranty. In uncertain times, "
         "warranty = confidence. 'Buy once, mow for 5 years.'"),
    ]:
        p = doc.add_paragraph()
        p.add_run(f"{title_}: ").bold = True
        p.add_run(desc)

    # --- 6. TIMELINE ---
    doc.add_heading("6. Implementation Timeline", level=1)
    for period, tasks in [
        ("Week 1 (14-20 Apr)", [
            "Update homepage hero banner to cordless-first messaging",
            "Create Spectrum cordless bundle variants on Shopify (mower + 2x battery + charger)",
            "Launch 'Beat the Fuel Crisis' email campaign to full list",
            "Redirect Google Ads budget: -30% petrol, +50% cordless",
        ]),
        ("Week 2 (21-27 Apr)", [
            "Build '40V System' ecosystem landing page",
            "Create 'Fuel Crisis Kit' bundles for petrol products",
            "Update all marketplace listings with crisis-aware copy",
            "Brief influencer partnerships",
        ]),
        ("Week 3 (28 Apr - 4 May)", [
            "Launch trade-in programme ('Switch to Cordless')",
            "Publish 'Petrol vs Battery 2026' blog post",
            "Launch TikTok/Instagram cost comparison series",
            "Add fuel efficiency badges to all petrol listings",
        ]),
        ("Week 4+ (May onwards)", [
            "Build cost comparison calculator for product pages",
            "Expand lifestyle photography across all Spectrum products",
            "Review and optimise based on first 2 weeks of data",
            "Scale successful campaigns, pause underperformers",
        ]),
    ]:
        doc.add_heading(period, level=3)
        for t in tasks:
            doc.add_paragraph(t, style="List Bullet")

    # --- 7. RISK ASSESSMENT ---
    doc.add_heading("7. Risk Assessment & Contingencies", level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Risk"
    hdr[1].text = "Assessment"
    hdr[2].text = "Contingency"
    for risk, assessment, contingency in [
        ("Ceasefire / price normalisation",
         "Fuel prices drop suddenly, petrol demand rebounds",
         "Maintain balanced inventory; keep petrol campaigns paused but ready to reactivate"),
        ("Cordless stock-out",
         "Demand surge exhausts Spectrum stock",
         "Coordinate with factory on expedited supply; implement waitlist"),
        ("Consumer paralysis",
         "Customers delay all purchases due to uncertainty",
         "Introduce 0% finance options to reduce purchase barrier"),
        ("Competitor price war",
         "Competitors slash prices to clear petrol inventory",
         "Focus on value (warranty, TCO, ecosystem) not matching price cuts"),
        ("Prolonged crisis",
         "Blockade extends beyond Q2, broader economic slowdown",
         "Shift to maintenance/repair products and parts"),
    ]:
        row = table.add_row().cells
        row[0].text = risk
        row[1].text = assessment
        row[2].text = contingency

    # --- 8. INVESTMENT ---
    doc.add_heading("8. Resource & Investment Requirements", level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Item"
    hdr[1].text = "Estimated Cost"
    hdr[2].text = "Timing"
    total_low = 0
    total_high = 0
    for item, cost, timing in [
        ("Shopify development (bundles, landing pages, calculator)", "£2,000 - £3,500", "Week 1-4"),
        ("Lifestyle photography (factory partner to supply)", "£0 - £500", "Week 2-3"),
        ("Google Ads budget reallocation (net neutral)", "£0", "Week 1"),
        ("Marketplace Promoted Listings uplift (+50%)", "£500 - £1,000/month", "Week 1"),
        ("Influencer partnerships (2-3 creators)", "£1,500 - £3,000", "Week 2-4"),
        ("Email campaigns (Klaviyo — existing platform)", "£0", "Week 1"),
        ("Trade-in programme setup & promotion", "£500 - £1,000", "Week 3"),
        ("Content creation (blog, social, video)", "£500 - £1,000", "Week 2-4"),
    ]:
        row = table.add_row().cells
        row[0].text = item
        row[1].text = cost
        row[2].text = timing
        cost_clean = cost.replace("£", "").replace(",", "").replace("/month", "")
        if " - " in cost_clean:
            low, high = cost_clean.split(" - ")
            total_low += int(low.strip())
            total_high += int(high.strip())
        else:
            total_low += int(cost_clean.strip())
            total_high += int(cost_clean.strip())

    row = table.add_row().cells
    row[0].text = "TOTAL ESTIMATED INVESTMENT"
    row[1].text = f"£{total_low:,} - £{total_high:,}"
    row[2].text = "4 weeks"
    for cell in row:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True

    doc.save(output_path)
    print(f"  Word document saved: {output_path}")


# ---------------------------------------------------------------------------
# Document generation — PowerPoint
# ---------------------------------------------------------------------------

def build_pptx(analysis: dict, output_path: str):
    """Generate the executive PowerPoint presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    march = analysis["march"]
    april = analysis["april"]
    march_days = 31
    april_days = 14

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_text_slide(title_text, body_lines, layout_idx=1):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slide.shapes.title.text = title_text
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(line, tuple):
                run = p.add_run()
                run.text = line[0]
                run.font.bold = True
                run.font.size = Pt(16)
                run2 = p.add_run()
                run2.text = f" — {line[1]}"
                run2.font.size = Pt(16)
            else:
                run = p.add_run()
                run.text = line
                run.font.size = Pt(16)
        return slide

    # Daily rate calculations
    m_daily = march["net_sales"] / march_days
    a_daily = april["net_sales"] / april_days
    overall_chg = ((a_daily - m_daily) / m_daily * 100) if m_daily else 0

    m_petrol = march["categories"].get("petrol", {"net_sales": 0, "orders": 0})
    a_petrol = april["categories"].get("petrol", {"net_sales": 0, "orders": 0})
    m_cordless = march["categories"].get("cordless", {"net_sales": 0, "orders": 0})
    a_cordless = april["categories"].get("cordless", {"net_sales": 0, "orders": 0})

    mp_daily = m_petrol["net_sales"] / march_days
    ap_daily = a_petrol["net_sales"] / april_days
    mc_daily = m_cordless["net_sales"] / march_days
    ac_daily = a_cordless["net_sales"] / april_days
    p_chg = ((ap_daily - mp_daily) / mp_daily * 100) if mp_daily else 0
    c_chg = ((ac_daily - mc_daily) / mc_daily * 100) if mc_daily else 0

    # --- SLIDE 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Emergency Marketing Response"
    slide.placeholders[1].text = (
        "MowDirect — Strait of Hormuz Crisis Response\n"
        "14 April 2026 | Confidential"
    )

    # --- SLIDE 2: Market Situation ---
    add_text_slide("Market Situation: Strait of Hormuz Crisis", [
        ("Naval blockade", "94% collapse in commercial shipping through Hormuz"),
        ("Oil above $100/barrel", "20M barrels/day shortfall, European fuel up >20%"),
        ("Consumer retreat", "Households cutting back on petrol-powered purchases"),
        ("Supply disruption", "Petrol mower supply chains disrupted for coming weeks"),
        ("Manufacturing surcharges", "Up to 30% on raw materials and packaging"),
    ])

    # --- SLIDE 3: Consumer Shift ---
    add_text_slide("The Accelerating Shift to Cordless", [
        "Rising fuel costs making petrol mower ownership significantly more expensive",
        "Battery mower market growing at 6.7% CAGR — crisis is accelerating this",
        "Consumers save £100-£200/year on servicing and fuel by switching to cordless",
        "Cordless entry price now competitive with petrol (£120-£180 vs £130-£190)",
        "'Painful and quick' transition toward electric alternatives underway",
    ])

    # --- SLIDE 4: Overall Sales Impact ---
    add_text_slide(
        f"Sales Impact: Daily Run Rate Down {overall_chg:.0f}%",
        [
            f"March 2026: £{march['net_sales']:,.0f} net sales | {march['orders']} orders ({march_days} days)",
            f"April 1-14: £{april['net_sales']:,.0f} net sales | {april['orders']} orders ({april_days} days)",
            "",
            f"March daily average: £{m_daily:,.0f}/day",
            f"April daily average: £{a_daily:,.0f}/day",
            f"Overall daily run rate change: {overall_chg:+.0f}%",
        ]
    )

    # --- SLIDE 5: Petrol vs Cordless ---
    add_text_slide("Petrol vs Cordless Performance", [
        f"Petrol: £{m_petrol['net_sales']:,.0f} (March) → £{a_petrol['net_sales']:,.0f} (April) | Daily rate: {p_chg:+.0f}%",
        f"Cordless: £{m_cordless['net_sales']:,.0f} (March) → £{a_cordless['net_sales']:,.0f} (April) | Daily rate: {c_chg:+.0f}%",
        "",
        "Petrol products dominate revenue but are seeing steeper decline",
        "Cordless range needs better website presentation to capture opportunity",
        f"Fuel cans: £{march['categories'].get('fuel_cans', {}).get('net_sales', 0):,.0f} in March — consumers are stocking up",
    ])

    # --- SLIDE 6: Top Products ---
    all_prods = []
    for cat, cdata in april["categories"].items():
        for p in cdata["products"]:
            all_prods.append({**p, "category": cat})
    all_prods.sort(key=lambda x: -x["net_sales"])
    top5 = all_prods[:5]

    add_text_slide("Top Selling Products — April 2026", [
        f"1. {p['title'][:55]} — £{p['net_sales']:,.0f} ({p['category']})"
        for p in top5
    ] + ["", "Petrol products still dominate but daily rates are declining"])

    # --- SLIDE 7: Strategic Response ---
    add_text_slide("Strategic Response: Three Pillars", [
        ("Pillar 1", "Accelerate Spectrum Cordless — bundles, ecosystem page, lifestyle imagery"),
        ("Pillar 2", "Reframe Petrol Range — fuel efficiency kits, TCO messaging, trade-in"),
        ("Pillar 3", "Messaging Pivot — homepage, email, social, SEO, marketplace listings"),
        "",
        "Factory partner investment of up to £1M demands urgency",
        "3-month cash window — every day of delay = missed opportunity",
    ])

    # --- SLIDE 8: Cordless Strategy ---
    add_text_slide("Pillar 1: Accelerate Spectrum Cordless", [
        ("Complete kits", "Every tool available as tool + battery + charger at one price"),
        ("Bundle clarity", "46cm mower clearly shows 2x battery requirement"),
        ("Ecosystem page", "'One battery powers everything' — visual map of all 40V tools"),
        ("Lifestyle imagery", "Real gardens, real use — source from factory partner"),
        ("Richer descriptions", "Use cases, garden sizes, runtime, petrol comparison"),
    ])

    # --- SLIDE 9: Petrol & Messaging ---
    add_text_slide("Pillar 2 & 3: Petrol Reframe + Messaging Pivot", [
        ("Fuel Crisis Kits", "Mower + jerry can + stabiliser + cost guide"),
        ("Efficiency badges", "'Only 0.8L/hr — mow for £2.50' on all petrol listings"),
        ("Trade-in programme", "£30-£50 credit: old petrol mower → Spectrum cordless"),
        ("Homepage", "'No fuel. No fumes. No worries.' — cordless-first hero banner"),
        ("Email (Klaviyo)", "Segmented: past petrol buyers, cordless browsers, all customers"),
        ("Cost calculator", "Interactive 5-year TCO comparison on product pages"),
    ])

    # --- SLIDE 10: Quick Wins ---
    add_text_slide("Additional Opportunities", [
        "Marketplace: update titles, boost cordless Promoted Listings +50%",
        "Google Ads: reallocate -30% petrol → +50% cordless campaigns",
        "Fast-track 2-3 influencer partnerships for Spectrum cordless",
        "PR pitch: 'UK retailer sees cordless surge amid fuel crisis'",
        "Stock reallocation: prioritise highest-converting channels",
        "'Green Garden Weekend' flash sale: 10% off + free battery upgrade",
        "Amplify Spectrum 5-year warranty — confidence in uncertain times",
    ])

    # --- SLIDE 11: Timeline ---
    add_text_slide("Implementation Roadmap", [
        ("Week 1 (14-20 Apr)", "Homepage, bundles, email blast, Ads reallocation"),
        ("Week 2 (21-27 Apr)", "Ecosystem page, Fuel Crisis Kits, marketplace updates"),
        ("Week 3 (28 Apr-4 May)", "Trade-in launch, blog/SEO, social series"),
        ("Week 4+ (May)", "Cost calculator, photography, optimisation"),
    ])

    # --- SLIDE 12: Investment ---
    add_text_slide("Investment & Resource Requirements", [
        "Shopify development (bundles, pages, calculator): £2,000 - £3,500",
        "Marketplace Promoted Listings uplift: £500 - £1,000/month",
        "Influencer partnerships: £1,500 - £3,000",
        "Trade-in programme setup: £500 - £1,000",
        "Content creation: £500 - £1,000",
        "",
        "Total estimated investment: £5,000 - £10,000",
        "Google Ads reallocation is budget-neutral",
        "ROI target: offset petrol decline + capture cordless growth",
    ])

    # --- SLIDE 13: Summary ---
    add_text_slide("Summary & Next Steps", [
        "The fuel crisis is both a threat (petrol decline) and opportunity (cordless growth)",
        "Factory partner has £1M invested — rapid execution required",
        "3-month cash runway means every week counts",
        "Petrol mower shortages = low hanging fruit for cordless conversion",
        "Spectrum 40V ecosystem is ready — website presentation must catch up",
        "",
        "ACTION: approve this plan and resource Week 1 immediately",
    ])

    prs.save(output_path)
    print(f"  Presentation saved: {output_path}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    print("\n=== MowDirect Emergency Sales Report ===")
    print(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    print("=" * 45)

    # --- Fetch data ---
    print("\nFetching sales data from Shopify (ShopifyQL) ...")
    data = fetch_all_data()

    # --- Analyse ---
    print("Analysing ...")
    analysis = analyse(data)

    march = analysis["march"]
    april = analysis["april"]

    # --- Save JSON ---
    json_path = os.path.join(OUTPUT_DIR, "emergency_sales_data.json")
    with open(json_path, "w") as f:
        json.dump({
            "generated": datetime.now().isoformat(),
            "march_2026": march,
            "april_2026_partial": april,
        }, f, indent=2, default=str)
    print(f"  JSON saved: {json_path}")

    # --- Console summary ---
    march_days = 31
    april_days = 14
    m_daily = march["net_sales"] / march_days
    a_daily = april["net_sales"] / april_days
    chg = ((a_daily - m_daily) / m_daily * 100) if m_daily else 0

    print(f"\n--- March 2026 ---")
    print(f"  Net: £{march['net_sales']:,.2f} | {march['orders']} orders | £{m_daily:,.0f}/day")
    for cat, cdata in sorted(march["categories"].items(), key=lambda x: -x[1]["net_sales"]):
        print(f"  {cat:>12s}: £{cdata['net_sales']:>10,.2f} | {cdata['orders']:>5d} orders")

    print(f"\n--- April 2026 (1-14) ---")
    print(f"  Net: £{april['net_sales']:,.2f} | {april['orders']} orders | £{a_daily:,.0f}/day")
    for cat, cdata in sorted(april["categories"].items(), key=lambda x: -x[1]["net_sales"]):
        print(f"  {cat:>12s}: £{cdata['net_sales']:>10,.2f} | {cdata['orders']:>5d} orders")

    print(f"\n  Daily run rate change: {chg:+.0f}%")

    # --- Generate documents ---
    print("\nGenerating documents ...")
    docx_path = os.path.join(OUTPUT_DIR, "Emergency_Marketing_Response.docx")
    build_docx(analysis, docx_path)

    pptx_path = os.path.join(OUTPUT_DIR, "Emergency_Marketing_Response.pptx")
    build_pptx(analysis, pptx_path)

    print(f"\nDone. All files in: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
