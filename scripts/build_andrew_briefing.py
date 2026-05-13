"""
Build the Monday briefing deck for Andrew on the marketing spend rebuttal +
forward plan, based on the MowDirect-rebranded PPT template.

Strips the 36 Iowa sample slides, then constructs new slides using the
template's layouts so all master decorations (logo, footer, fonts) carry
through. Big-number slides use manually positioned text boxes so the stats
read as stats, not bullets.

Story arc (12 slides):
   1. Title
   2. Three things in this deck
   3. Section 01 — The £6k explained
   4. Where the £6k actually came from (card transactions + declines)
   5. Section 02 — April → This week, where we actually are
   6. ROAS comparison: 6.66× → 10.20× (+53%)
   7. Cash maths: same spend, more sales
   8. What's inside the 15.4% line
   9. Section 03 — The plan from here
  10. Three levers
  11. The ask
  12. Close

Output: reports/andrew_briefing_2026-05-11.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TEMPLATE = Path("/Users/waynetheisinger/Desktop/PPT-Template-Standard-2025.pptx")
OUTPUT   = Path("reports/andrew_briefing_2026-05-11.pptx")

# Brand-consistent palette pulled from the template's gold accent + black.
GOLD  = RGBColor(0xFF, 0xCD, 0x00)   # template gold
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY  = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY = RGBColor(0xAA, 0xAA, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xC0, 0x39, 0x2B)

# Layout indices (validated by inspect_template.py)
L_TITLE_BLACK   = 0
L_TITLE_GOLD    = 1
L_SECTION_GOLD  = 4
L_SECTION_BLACK = 5
L_TITLE_ONLY    = 7
L_BULLET        = 8
L_3x1_GRID      = 13
L_3COL_TEXT     = 11
L_3COL_ICON     = 18
L_CLOSING       = 25


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def strip_all_existing_slides(prs):
    """Remove every slide from the deck, leaving the master + layouts intact."""
    sldIdLst = prs.slides._sldIdLst
    rels = prs.part.rels
    for sldId in list(sldIdLst):
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        sldIdLst.remove(sldId)
        if rId in rels:
            rels.pop(rId)


def set_text(placeholder, text, *, size=None, bold=None, color=None, align=None):
    """Write text into a placeholder, preserving style where not overridden."""
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    if size  is not None: run.font.size  = Pt(size)
    if bold  is not None: run.font.bold  = bold
    if color is not None: run.font.color.rgb = color


def add_textbox(slide, *, left, top, width, height, text,
                size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def find_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def add_slide(prs, layout_idx):
    layout = prs.slide_masters[0].slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def fill_title_body(slide, title, body_lines):
    title_ph = find_placeholder(slide, 0)
    body_ph  = find_placeholder(slide, 1)
    if title_ph is not None:
        set_text(title_ph, title)
    if body_ph is not None:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = add_slide(prs, L_TITLE_BLACK)
    title_ph = find_placeholder(slide, 0)
    if title_ph:
        set_text(title_ph, "Marketing Spend Review")
    body_ph = find_placeholder(slide, 1)
    if body_ph:
        set_text(body_ph, "BRIEFING FOR ANDREW GUNN — 11 MAY 2026")
    sub_ph = find_placeholder(slide, 21)
    if sub_ph:
        set_text(sub_ph, "Wayne Theisinger · MowDirect")
    return slide


def slide_toc(prs):
    slide = add_slide(prs, L_BULLET)
    fill_title_body(
        slide,
        "Three things in this deck",
        [
            "Where the £6k on the card actually came from",
            "April → This week: spend held, ROAS up 53%",
            "The plan from here — three levers, none requiring more ad spend",
        ],
    )
    return slide


def slide_section(prs, number, title, subtitle=None):
    slide = add_slide(prs, L_SECTION_GOLD)
    title_ph = find_placeholder(slide, 0)
    if title_ph:
        set_text(title_ph, f"{number}. {title}")
    body_ph = find_placeholder(slide, 1)
    if body_ph and subtitle:
        set_text(body_ph, subtitle)
    return slide


# ---------------------------------------------------------------------------
# Section 01 — The £6k explained
# ---------------------------------------------------------------------------

def slide_6k_receipts(prs):
    """The actual card transactions. Three blocks: this week / 14-day / declines."""
    slide = add_slide(prs, L_TITLE_ONLY)
    title_ph = find_placeholder(slide, 0)
    if title_ph:
        set_text(title_ph, "Where the £6k actually came from")

    slide_w = prs.slide_width
    margin  = Inches(0.5)
    gutter  = Inches(0.25)
    col_w   = (slide_w - 2 * margin - gutter) // 2
    top     = Inches(1.9)

    # LEFT block — this week
    left = margin
    add_textbox(slide, left=left, top=top, width=col_w, height=Inches(0.4),
                text="MON 4 – SUN 10 MAY (\"this week\")",
                size=11, bold=True, color=GOLD,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left=left, top=top + Inches(0.45), width=col_w, height=Inches(1.4),
                text="£2,500",
                size=80, bold=True, color=BLACK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left=left, top=top + Inches(1.9), width=col_w, height=Inches(0.4),
                text="across 4 card charges",
                size=14, color=GREY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Daily breakdown under the £2,500
    lines = [
        "  Tue 6 May    £500",
        "  Tue 6 May    £1,000",
        "  Wed 7 May    £500",
        "  Thu 8 May    £500",
    ]
    add_textbox(slide, left=left + Inches(0.3), top=top + Inches(2.4),
                width=col_w - Inches(0.6), height=Inches(1.4),
                text="\n".join(lines),
                size=14, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # RIGHT block — last 14 days
    left = margin + col_w + gutter
    add_textbox(slide, left=left, top=top, width=col_w, height=Inches(0.4),
                text="LAST 14 DAYS (26 APR – 10 MAY)",
                size=11, bold=True, color=GOLD,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left=left, top=top + Inches(0.45), width=col_w, height=Inches(1.4),
                text="£5,791",
                size=80, bold=True, color=BLACK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left=left, top=top + Inches(1.9), width=col_w, height=Inches(0.4),
                text="across 10 charges (including April month-end)",
                size=14, color=GREY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Breakdown right column
    lines_r = [
        "  Apr 26–30      £2,500   (April's last week)",
        "  May 1           £291   (April month-end true-up)",
        "  May 2           £500",
        "  May 6 → 8      £2,500   (this week)",
    ]
    add_textbox(slide, left=left + Inches(0.2), top=top + Inches(2.4),
                width=col_w - Inches(0.4), height=Inches(1.4),
                text="\n".join(lines_r),
                size=14, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Footer takeaway
    add_textbox(slide,
                left=margin,
                top=prs.slide_height - Inches(1.4),
                width=slide_w - 2 * margin, height=Inches(0.55),
                text="The £6k is two weeks of card charges, not one — and starts in April.",
                size=20, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                left=margin,
                top=prs.slide_height - Inches(0.85),
                width=slide_w - 2 * margin, height=Inches(0.4),
                text="Plus 4 declined charges in May (Mon 3 · Tue 5 · Sat 9) — that is a Visa cash flow issue, not a marketing overspend.",
                size=12, color=RED, align=PP_ALIGN.CENTER)
    return slide


# ---------------------------------------------------------------------------
# Section 02 — April → This week
# ---------------------------------------------------------------------------

def slide_roas_jump(prs):
    """The headline improvement — ROAS jump from April."""
    slide = add_slide(prs, L_TITLE_ONLY)
    title_ph = find_placeholder(slide, 0)
    if title_ph:
        set_text(title_ph, "Every £1 spent is returning 53% more sales")

    slide_w = prs.slide_width
    margin  = Inches(0.5)
    gutter  = Inches(0.4)
    col_w   = (slide_w - 2 * margin - 2 * gutter) // 3
    top     = Inches(2.0)

    blocks = [
        ("APRIL", "6.66×", "ROAS\nthe baseline", GREY),
        ("", "→", "", LIGHT_GREY),
        ("THIS WEEK", "10.20×", "ROAS\nMon 4 – Sun 10 May", BLACK),
    ]
    for i, (tag, big, body, color) in enumerate(blocks):
        left = margin + i * (col_w + gutter)
        add_textbox(slide, left=left, top=top, width=col_w, height=Inches(0.4),
                    text=tag, size=12, bold=True, color=GOLD,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=left, top=top + Inches(0.5), width=col_w, height=Inches(2.0),
                    text=big, size=110 if i != 1 else 80, bold=True, color=color,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=left, top=top + Inches(2.6), width=col_w, height=Inches(0.9),
                    text=body, size=14, color=GREY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Big delta callout
    add_textbox(slide,
                left=margin, top=Inches(5.4),
                width=slide_w - 2 * margin, height=Inches(0.65),
                text="+53% improvement in efficiency",
                size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                left=margin, top=Inches(6.05),
                width=slide_w - 2 * margin, height=Inches(0.5),
                text="Same daily spend (~£408/day). 53% more sales per £1.",
                size=14, color=GREY, align=PP_ALIGN.CENTER)
    return slide


def slide_cash_maths(prs):
    """The actual cash numbers — spend and resulting tracked sales."""
    slide = add_slide(prs, L_TITLE_ONLY)
    title_ph = find_placeholder(slide, 0)
    if title_ph:
        set_text(title_ph, "What that means in actual cash")

    slide_w = prs.slide_width
    margin  = Inches(0.5)
    gutter  = Inches(0.3)
    col_w   = (slide_w - 2 * margin - gutter) // 2
    top     = Inches(1.9)

    blocks = [
        ("APRIL — 30 DAYS",
         "£12,282", "Google Ads spend",
         "£81,788", "Tracked sales attributed",
         "6.66×", "ROAS"),
        ("THIS WEEK — 7 DAYS",
         "£2,856", "Google Ads spend",
         "£29,146", "Tracked sales attributed",
         "10.20×", "ROAS"),
    ]

    for i, (header, big1, lbl1, big2, lbl2, big3, lbl3) in enumerate(blocks):
        left = margin + i * (col_w + gutter)
        # Section header
        add_textbox(slide, left=left, top=top, width=col_w, height=Inches(0.5),
                    text=header, size=14, bold=True, color=GOLD,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Three rows
        y = top + Inches(0.7)
        for big, lbl in [(big1, lbl1), (big2, lbl2), (big3, lbl3)]:
            add_textbox(slide, left=left, top=y, width=col_w, height=Inches(0.7),
                        text=big, size=44, bold=True, color=BLACK,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_textbox(slide, left=left, top=y + Inches(0.75), width=col_w, height=Inches(0.35),
                        text=lbl, size=11, color=GREY,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            y = y + Inches(1.25)

    # Footer takeaway
    add_textbox(slide,
                left=margin,
                top=prs.slide_height - Inches(1.0),
                width=slide_w - 2 * margin, height=Inches(0.55),
                text="ClickSlice running tighter, not destroyed.",
                size=20, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                left=margin,
                top=prs.slide_height - Inches(0.5),
                width=slide_w - 2 * margin, height=Inches(0.4),
                text="The data is right. The cash position is a separate question — see Section 01.",
                size=12, color=GREY, align=PP_ALIGN.CENTER)
    return slide


def slide_15pct_breakdown(prs):
    """Three-column breakdown of what Andrew sees as the marketing line."""
    slide = add_slide(prs, L_TITLE_ONLY)
    title_ph = find_placeholder(slide, 0)
    if title_ph:
        set_text(title_ph, "What's inside the marketing line — April")

    cols = [
        ("3.5%",  "Paid ad spend",         "FLEXIBLE",   "Google Ads. This is marketing.\nWell under the 10% allowance."),
        ("9.0%",  "Marketplace commissions","STRUCTURAL", "Amazon, eBay, B&Q referral fees.\nCost of being on those channels, not advertising."),
        ("2.9%",  "FBA fulfilment",        "OPERATIONS", "Pick · pack · ship · storage.\nOutward delivery, paid to Amazon."),
    ]

    slide_w = prs.slide_width
    margin  = Inches(0.5)
    gutter  = Inches(0.2)
    col_w   = (slide_w - 2 * margin - 2 * gutter) // 3
    top     = Inches(2.0)
    h_big   = Inches(1.3)
    h_label = Inches(0.5)
    h_tag   = Inches(0.4)
    h_body  = Inches(2.0)

    for i, (pct, label, tag, body) in enumerate(cols):
        left = margin + i * (col_w + gutter)
        add_textbox(slide, left=left, top=top, width=col_w, height=h_big,
                    text=pct, size=72, bold=True,
                    color=GOLD if i == 0 else BLACK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=left, top=top + h_big, width=col_w, height=h_label,
                    text=label, size=18, bold=True, color=BLACK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=left, top=top + h_big + h_label, width=col_w, height=h_tag,
                    text=tag, size=11, bold=True,
                    color=GOLD if i == 0 else GREY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=left + Inches(0.1),
                    top=top + h_big + h_label + h_tag,
                    width=col_w - Inches(0.2), height=h_body,
                    text=body, size=12, color=GREY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    add_textbox(slide,
                left=margin,
                top=prs.slide_height - Inches(1.0),
                width=slide_w - 2 * margin, height=Inches(0.55),
                text="Only the gold column is marketing. The other 11.9% is structural.",
                size=18, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                left=margin,
                top=prs.slide_height - Inches(0.5),
                width=slide_w - 2 * margin, height=Inches(0.4),
                text="Treat Amazon, eBay, B&Q fees the way you treat outward delivery — operations, not marketing.",
                size=12, color=GREY, align=PP_ALIGN.CENTER)
    return slide


# ---------------------------------------------------------------------------
# Section 03 — The plan from here
# ---------------------------------------------------------------------------

def slide_three_levers(prs):
    slide = add_slide(prs, L_TITLE_ONLY)
    title_ph = find_placeholder(slide, 0)
    if title_ph:
        set_text(title_ph, "Three levers — none requiring more ad spend")

    levers = [
        ("SALESFIRE",
         "Re-engage the 13,600 customers we've already paid to acquire",
         "Email + SMS + onsite personalisation in one platform. Signed in April; getting operational this month. Near-zero marginal cost per re-engagement."),
        ("PER-SKU CONTRIBUTION",
         "Net contribution after referrals + FBA + ad spend",
         "New tab in the monthly report. Surfaces SKUs that are structurally underwater on each marketplace. Drives price discipline at the SKU level."),
        ("CHANNEL MIX",
         "Bias revenue from marketplaces to Shopify direct",
         "Every £1 moved direct saves ~10p of structural channel cost (~12% on marketplaces vs ~2% direct). Spectrum-first paid + Awin affiliate push this."),
    ]

    slide_w = prs.slide_width
    margin  = Inches(0.5)
    gutter  = Inches(0.2)
    col_w   = (slide_w - 2 * margin - 2 * gutter) // 3
    top     = Inches(2.0)

    for i, (tag, headline, body) in enumerate(levers):
        left = margin + i * (col_w + gutter)
        add_textbox(slide, left=left, top=top, width=col_w, height=Inches(0.45),
                    text=tag, size=12, bold=True, color=GOLD,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=left, top=top + Inches(0.45),
                    width=col_w, height=Inches(1.5),
                    text=headline, size=18, bold=True, color=BLACK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        add_textbox(slide, left=left + Inches(0.1), top=top + Inches(2.0),
                    width=col_w - Inches(0.2), height=Inches(2.5),
                    text=body, size=12, color=GREY,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    return slide


def slide_the_ask(prs):
    slide = add_slide(prs, L_BULLET)
    fill_title_body(
        slide,
        "What I need from you",
        [
            "Agree the marketing budget is measured against paid ad spend (currently 3.5% of gross)",
            "Hold the 10% allowance against that bucket only",
            "Treat marketplace commissions and FBA fulfilment as separate operational lines",
            "Top up the Visa — 4 declines in May is the cash flow signal, not the spend",
        ],
    )
    return slide


def slide_close(prs):
    slide = add_slide(prs, L_CLOSING)
    title_ph = find_placeholder(slide, 0)
    if title_ph:
        set_text(title_ph, "Questions?")
    body_ph = find_placeholder(slide, 1)
    if body_ph:
        set_text(body_ph, "Wayne Theisinger")
    sub_ph = find_placeholder(slide, 21)
    if sub_ph:
        set_text(sub_ph, "wayne@mowdirect.co.uk")
    return slide


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation(str(TEMPLATE))
    strip_all_existing_slides(prs)

    slide_title(prs)
    slide_toc(prs)

    slide_section(prs, "01", "The £6k explained",
                  "What actually hit the card")
    slide_6k_receipts(prs)

    slide_section(prs, "02", "April → This week",
                  "Spend held — efficiency jumped 53%")
    slide_roas_jump(prs)
    slide_cash_maths(prs)
    slide_15pct_breakdown(prs)

    slide_section(prs, "03", "The plan from here",
                  "Operate smarter — not spend more")
    slide_three_levers(prs)
    slide_the_ask(prs)

    slide_close(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}  ({OUTPUT.stat().st_size:,} bytes)")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    build()
